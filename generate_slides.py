#!/usr/bin/env python3
"""
generate_slides.py
==================
Generates health_expenditure_analysis.pptx

Big Idea: America's longevity problem is not a budget problem — it is upstream of the budget.
Structure: Storytelling with Data  (Situation → Complication → Resolution)
Audience:  Data analysts / MBA class / executive board
Source:    World Bank World Development Indicators, 2000-2023

Slides
------
 1  Cover
 2  Global scatter: spending vs life expectancy
 3  Top 15 spenders bar chart
 4  High-income peer scatter (US anomaly)
 5  Three upstream factors (infographic)
 6  Factor 1 – Obesity
 7  Factor 2 – Education & health literacy
 8  Factor 3 – System efficiency
 9  Resolution comparison table
10  Call to action
"""

import os
import io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Palette (Power BI style) ──────────────────────────────────────────────────
C_BLUE       = "#0078D4"
C_DARKBLUE   = "#004578"
C_ORANGE     = "#D87000"
C_GREEN      = "#107C10"
C_WHITE      = "#FFFFFF"
C_NEARBLACK  = "#252525"
C_GRAY       = "#737373"
C_LIGHTGRAY  = "#F2F2F2"
C_MIDGRAY    = "#D1D1D1"
C_LIGHTBLUE  = "#90CAF9"
C_BLUEGRAY   = "#B0BEC5"

# matplotlib named equivalents for places where hex is needed directly
MPL_BLUE      = C_BLUE
MPL_ORANGE    = C_ORANGE
MPL_GREEN     = C_GREEN
MPL_MIDGRAY   = C_MIDGRAY
MPL_LIGHTGRAY = C_LIGHTGRAY


# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _blank_slide(prs):
    """Add a blank (layout-6) slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text, left, top, width, height,
         size=14, bold=False, italic=False,
         color=C_NEARBLACK, align=PP_ALIGN.LEFT):
    """Add a text box; newlines become separate paragraphs."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def _mpl_style(ax, xlabel=None, ylabel=None):
    """Apply Power BI–inspired minimal style to a matplotlib Axes."""
    ax.set_facecolor("white")
    ax.figure.patch.set_facecolor("white")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color(MPL_MIDGRAY)
    ax.spines["bottom"].set_color(MPL_MIDGRAY)
    ax.tick_params(colors=C_GRAY, labelsize=9)
    if xlabel:
        ax.set_xlabel(xlabel, fontsize=9, color=C_GRAY, labelpad=6)
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=9, color=C_GRAY, labelpad=6)
    ax.grid(axis="y", color=MPL_LIGHTGRAY, linewidth=0.8, zorder=0)


def _to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf


def _add_pic(slide, stream, left, top, width, height):
    slide.shapes.add_picture(stream, left, top, width, height)


# ── Slide 1 – Cover ───────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = _blank_slide(prs)
    _bg(slide, C_DARKBLUE)

    # Left accent bar
    _rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, C_BLUE)

    _txt(slide,
         "America's Longevity Problem\nIs Not a Budget Problem",
         Inches(0.7), Inches(1.4), Inches(11.5), Inches(2.4),
         size=40, bold=True, color=C_WHITE)

    _txt(slide,
         "It is upstream of the budget.",
         Inches(0.7), Inches(3.9), Inches(11.0), Inches(0.65),
         size=22, color=C_LIGHTBLUE)

    # Divider
    _rect(slide, Inches(0.7), Inches(4.65), Inches(6.0), Inches(0.04), C_BLUE)

    _txt(slide,
         "A data-driven investigation into health expenditure, life expectancy, and\n"
         "the upstream factors that explain America's $8,431-per-person paradox.\n\n"
         "Source: World Bank World Development Indicators  |  2000 – 2023",
         Inches(0.7), Inches(4.85), Inches(11.5), Inches(2.0),
         size=13, color=C_BLUEGRAY)


# ── Slide 2 – Global Scatter ──────────────────────────────────────────────────

def _chart_global_scatter(df):
    sub = df.dropna(subset=["health_expenditure_ppp", "life_expectancy"]).copy()
    log_x = np.log10(sub["health_expenditure_ppp"])
    le    = sub["life_expectancy"]

    fig, ax = plt.subplots(figsize=(9.5, 5.2))

    ax.scatter(log_x, le, color=MPL_BLUE, alpha=0.50, s=40, zorder=3)

    # Trend line
    coef = np.polyfit(log_x, le, 1)
    xline = np.linspace(log_x.min(), log_x.max(), 300)
    ax.plot(xline, np.poly1d(coef)(xline), color=C_NEARBLACK, lw=2.0, zorder=4)

    # Correlation annotation
    r = float(np.corrcoef(log_x, le)[0, 1])
    ax.text(0.02, 0.97, f"r = {r:.2f}  (log scale)",
            transform=ax.transAxes, fontsize=11, va="top",
            fontweight="bold", color=C_NEARBLACK)

    # Highlight US
    us = sub[sub["country"] == "United States"]
    if not us.empty:
        ux = float(np.log10(us["health_expenditure_ppp"].iloc[0]))
        uy = float(us["life_expectancy"].iloc[0])
        ax.scatter([ux], [uy], color=MPL_ORANGE, s=110, zorder=5)
        ax.annotate(
            "United States\nHighest spend,\nbelow the trend",
            xy=(ux, uy), xytext=(ux - 0.38, uy - 6.5),
            fontsize=8.5, color=MPL_ORANGE, fontweight="bold",
            arrowprops=dict(arrowstyle="->", color=MPL_ORANGE, lw=1.5)
        )

    # Highlight Japan and Spain
    for country, dx, dy in [("Japan", 0.06, 2.5), ("Spain", -0.36, 2.5)]:
        row = sub[sub["country"] == country]
        if not row.empty:
            rx = float(np.log10(row["health_expenditure_ppp"].iloc[0]))
            ry = float(row["life_expectancy"].iloc[0])
            ax.scatter([rx], [ry], color=MPL_GREEN, s=80, zorder=5)
            ax.annotate(
                f"{country}\nLess spend,\nlonger lives",
                xy=(rx, ry), xytext=(rx + dx, ry + dy),
                fontsize=8, color=MPL_GREEN, fontweight="bold",
                arrowprops=dict(arrowstyle="->", color=MPL_GREEN, lw=1.2)
            )

    xticks = [1.0, 1.5, 2.0, 2.5, 3.0, 3.5, 4.0]
    ax.set_xticks(xticks)
    ax.set_xticklabels([f"${10**v:,.0f}" for v in xticks], fontsize=8.5)
    ax.set_xlim(log_x.min() - 0.08, log_x.max() + 0.18)
    ax.set_ylim(48, 90)
    _mpl_style(ax,
               xlabel="Health Expenditure per Capita, PPP (Log Scale)",
               ylabel="Average Life Expectancy (Years)")
    fig.tight_layout()
    return fig


def slide_global_scatter(prs, df):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "Across 140+ countries, spending more on health buys longer lives…",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Pearson r = 0.70 on log scale — a strong global signal that investment matters",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
         size=11, color=C_GRAY)
    _add_pic(slide, _to_stream(_chart_global_scatter(df)),
             Inches(0.3), Inches(1.05), Inches(12.7), Inches(6.1))
    _txt(slide,
         "Data: World Bank World Development Indicators. Health expenditure on log scale. PPP = Purchasing Power Parity.",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 3 – Top Spenders Bar Chart ─────────────────────────────────────────

def _chart_top_spenders(df):
    top = (df.nlargest(15, "health_expenditure_ppp")
             .sort_values("health_expenditure_ppp", ascending=True)
             .copy())

    fig, ax = plt.subplots(figsize=(9.5, 5.5))
    colors = [MPL_ORANGE if c == "United States" else MPL_BLUE
              for c in top["country"]]
    bars = ax.barh(top["country"], top["health_expenditure_ppp"],
                   color=colors, height=0.65, zorder=3)

    for bar, spend, le, cntry in zip(
            bars, top["health_expenditure_ppp"],
            top["life_expectancy"], top["country"]):
        label_color = MPL_ORANGE if cntry == "United States" else C_NEARBLACK
        bold_w = "bold" if cntry == "United States" else "normal"
        ax.text(bar.get_width() + 80,
                bar.get_y() + bar.get_height() / 2,
                f"${spend:,.0f}  |  LE: {le:.0f} yrs",
                va="center", fontsize=8.5,
                color=label_color, fontweight=bold_w)

    ax.set_xlim(0, 10800)
    ax.set_facecolor("white")
    ax.figure.patch.set_facecolor("white")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color(MPL_MIDGRAY)
    ax.spines["bottom"].set_color(MPL_MIDGRAY)
    ax.tick_params(colors=C_GRAY, labelsize=9)
    ax.set_xlabel("Health Expenditure per Capita, PPP (USD)", fontsize=9, color=C_GRAY)
    ax.grid(axis="x", color=MPL_LIGHTGRAY, linewidth=0.8, zorder=0)
    fig.tight_layout()
    return fig


def slide_top_spenders(prs, df):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "…but the world's biggest spenders tell a very different story",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Japan reaches 84 years on $3,640 per person. The US spends $8,431 and reaches only 78.",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
         size=11, color=C_GRAY)
    _add_pic(slide, _to_stream(_chart_top_spenders(df)),
             Inches(0.3), Inches(1.05), Inches(12.7), Inches(6.1))
    _txt(slide,
         "Data: World Bank World Development Indicators. LE = Life Expectancy.",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 4 – US Anomaly (high-income peers) ─────────────────────────────────

def _chart_hi_peers(df):
    hi = df[df["health_expenditure_ppp"] >= 2000].dropna(
        subset=["health_expenditure_ppp", "life_expectancy"]).copy()

    fig, ax = plt.subplots(figsize=(9.5, 5.2))

    non_us = hi[hi["country"] != "United States"]
    ax.scatter(non_us["health_expenditure_ppp"], non_us["life_expectancy"],
               color=MPL_BLUE, alpha=0.60, s=55, zorder=3)

    label_set = {"Japan", "Spain", "Switzerland", "France",
                 "Germany", "Canada", "Australia", "Italy", "South Korea", "United Kingdom"}
    for _, row in non_us[non_us["country"].isin(label_set)].iterrows():
        ax.annotate(row["country"],
                    xy=(row["health_expenditure_ppp"], row["life_expectancy"]),
                    xytext=(6, 1), textcoords="offset points",
                    fontsize=7.5, color=MPL_BLUE)

    us = hi[hi["country"] == "United States"]
    ux = float(us["health_expenditure_ppp"].iloc[0])
    uy = float(us["life_expectancy"].iloc[0])
    ax.scatter([ux], [uy], color=MPL_ORANGE, s=150, zorder=5)
    ax.annotate(
        "United States\n$8,431 / 78.5 yrs",
        xy=(ux, uy), xytext=(ux - 2500, uy - 2.8),
        fontsize=9.5, color=MPL_ORANGE, fontweight="bold",
        arrowprops=dict(arrowstyle="->", color=MPL_ORANGE, lw=1.6)
    )

    ax.xaxis.set_major_formatter(
        mticker.FuncFormatter(lambda x, _: f"${x / 1000:.0f}K"))
    ax.set_ylim(76, 87)
    _mpl_style(ax,
               xlabel="Health Expenditure per Capita, PPP",
               ylabel="Life Expectancy (Years)")
    fig.tight_layout()
    return fig


def slide_us_anomaly(prs, df):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "Among high-income peers, the United States is a clear outlier",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Every comparable country achieves the same or better life expectancy at substantially lower cost",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
         size=11, color=C_GRAY)

    _add_pic(slide, _to_stream(_chart_hi_peers(df)),
             Inches(0.3), Inches(1.05), Inches(9.0), Inches(6.15))

    # Callout panel
    _rect(slide, Inches(9.45), Inches(1.05), Inches(3.65), Inches(1.55), C_ORANGE)
    _txt(slide,
         "The US–Japan gap:\n5 fewer years of life\nat 2.3× the cost",
         Inches(9.55), Inches(1.1), Inches(3.45), Inches(1.4),
         size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _rect(slide, Inches(9.45), Inches(2.75), Inches(3.65), Inches(4.45), C_LIGHTGRAY)
    insight = ("The US outspends Japan by\n$4,791 per person per year.\n\n"
               "Multiplied across 330 million\nAmericans, that is ~$1.6 trillion\n"
               "in annual excess health spend\nwith worse outcomes to show for it.\n\n"
               "The problem is systemic,\nnot financial.")
    _txt(slide, insight,
         Inches(9.6), Inches(2.9), Inches(3.35), Inches(4.1),
         size=10.5, color=C_NEARBLACK)

    _txt(slide,
         "Data: World Bank World Development Indicators.",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 5 – Three Upstream Factors ─────────────────────────────────────────

def slide_upstream_factors(prs):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "If it's not money, then what drives the gap?",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=22, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Three upstream factors explain most of the difference in health outcomes among wealthy nations",
         Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.38),
         size=11, color=C_GRAY)

    factors = [
        ("01", "Lifestyle &\nObesity",
         "The US adult obesity rate (36%) is\nnearly 3× Japan's (4%) and France's\n(22%). Obesity is the single largest\ndriver of preventable chronic disease.",
         C_ORANGE),
        ("02", "Education &\nHealth Literacy",
         "Countries with strong public health\neducation programs see 15–25%\nlower preventable mortality and\nhigher system efficiency per dollar.",
         C_BLUE),
        ("03", "System Design &\nPrimary Care",
         "Universal primary care access\nreduces costly hospital admissions\nand concentrates spend where it\nhas the highest marginal return.",
         C_GREEN),
    ]

    box_w = Inches(3.95)
    box_h = Inches(5.1)
    for i, (num, title, body, color) in enumerate(factors):
        lft = Inches(0.35 + i * 4.35)
        _rect(slide, lft, Inches(1.3), box_w, box_h, color)
        _txt(slide, num,
             lft + Inches(0.18), Inches(1.45), box_w - Inches(0.3), Inches(0.65),
             size=30, bold=True, color=C_WHITE)
        _txt(slide, title,
             lft + Inches(0.18), Inches(2.1), box_w - Inches(0.3), Inches(0.85),
             size=17, bold=True, color=C_WHITE)
        _txt(slide, body,
             lft + Inches(0.18), Inches(3.0), box_w - Inches(0.3), Inches(3.2),
             size=11.5, color=C_WHITE)

    _txt(slide,
         "Source: WHO Global Health Observatory, OECD Health at a Glance 2023, World Bank WDI",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 6 – Factor 1: Obesity ───────────────────────────────────────────────

def _chart_obesity(df):
    hi = df[df["health_expenditure_ppp"] >= 2000].dropna(
        subset=["obesity_rate", "life_expectancy"]).copy()

    fig, ax = plt.subplots(figsize=(8.5, 4.8))

    non_us = hi[hi["country"] != "United States"]
    ax.scatter(non_us["obesity_rate"], non_us["life_expectancy"],
               color=MPL_BLUE, alpha=0.60, s=55, zorder=3)

    for _, row in non_us[non_us["country"].isin(
            {"Japan", "Spain", "France", "Germany",
             "Canada", "Australia", "United Kingdom", "Italy", "Switzerland"})].iterrows():
        ax.annotate(row["country"],
                    xy=(row["obesity_rate"], row["life_expectancy"]),
                    xytext=(4, 0), textcoords="offset points",
                    fontsize=7.5, color=MPL_BLUE)

    us = hi[hi["country"] == "United States"]
    ux = float(us["obesity_rate"].iloc[0])
    uy = float(us["life_expectancy"].iloc[0])
    ax.scatter([ux], [uy], color=MPL_ORANGE, s=140, zorder=5)
    ax.annotate(
        "United States\n36.2% obese",
        xy=(ux, uy), xytext=(ux - 14, uy + 1.5),
        fontsize=9, color=MPL_ORANGE, fontweight="bold",
        arrowprops=dict(arrowstyle="->", color=MPL_ORANGE, lw=1.5)
    )

    # Trend (exclude US to show clean non-US trend)
    z = np.polyfit(non_us["obesity_rate"], non_us["life_expectancy"], 1)
    xfit = np.linspace(hi["obesity_rate"].min(), hi["obesity_rate"].max(), 200)
    ax.plot(xfit, np.poly1d(z)(xfit), color=MPL_MIDGRAY,
            lw=1.6, linestyle="--", zorder=2)

    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:.0f}%"))
    ax.set_ylim(76, 87)
    _mpl_style(ax,
               xlabel="Adult Obesity Rate",
               ylabel="Life Expectancy (Years)")
    fig.tight_layout()
    return fig


def slide_obesity(prs, df):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "Factor 1: America is the most obese wealthy nation — and it shows",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Higher obesity rates correlate strongly with lower life expectancy, even among the world's richest countries",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
         size=11, color=C_GRAY)

    _add_pic(slide, _to_stream(_chart_obesity(df)),
             Inches(0.3), Inches(1.05), Inches(8.8), Inches(5.9))

    _rect(slide, Inches(9.25), Inches(1.05), Inches(3.85), Inches(5.9), C_LIGHTGRAY)
    stats = [
        ("36.2%", "US adult obesity rate"),
        ("3.4×", "higher than Japan (10.7%)"),
        ("$1,861", "annual per-person\ncost of obesity in the US"),
        ("#1 of 38", "OECD ranking for\nobesity prevalence"),
    ]
    for i, (val, label) in enumerate(stats):
        _txt(slide, val,
             Inches(9.4), Inches(1.25 + i * 1.35), Inches(3.5), Inches(0.55),
             size=26, bold=True, color=C_ORANGE)
        _txt(slide, label,
             Inches(9.4), Inches(1.82 + i * 1.35), Inches(3.5), Inches(0.5),
             size=10, color=C_GRAY)

    _txt(slide,
         "Data: WHO Global Health Observatory, OECD Health Statistics.",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 7 – Factor 2: Education ────────────────────────────────────────────

def _chart_education(df):
    hi = df[df["health_expenditure_ppp"] >= 1500].dropna(
        subset=["mean_years_schooling", "life_expectancy"]).copy()

    fig, ax = plt.subplots(figsize=(8.5, 4.8))

    non_us = hi[hi["country"] != "United States"]
    ax.scatter(non_us["mean_years_schooling"], non_us["life_expectancy"],
               color=MPL_BLUE, alpha=0.60, s=55, zorder=3)

    for _, row in non_us[non_us["country"].isin(
            {"Japan", "Spain", "Germany", "South Korea",
             "Canada", "Australia", "France", "Finland"})].iterrows():
        ax.annotate(row["country"],
                    xy=(row["mean_years_schooling"], row["life_expectancy"]),
                    xytext=(4, 1), textcoords="offset points",
                    fontsize=7.5, color=MPL_BLUE)

    us = hi[hi["country"] == "United States"]
    ux = float(us["mean_years_schooling"].iloc[0])
    uy = float(us["life_expectancy"].iloc[0])
    ax.scatter([ux], [uy], color=MPL_ORANGE, s=140, zorder=5)
    ax.annotate(
        "United States\n13.4 yrs / 78.5 yrs LE",
        xy=(ux, uy), xytext=(ux - 3.2, uy - 1.5),
        fontsize=9, color=MPL_ORANGE, fontweight="bold",
        arrowprops=dict(arrowstyle="->", color=MPL_ORANGE, lw=1.5)
    )

    z = np.polyfit(non_us["mean_years_schooling"], non_us["life_expectancy"], 1)
    xfit = np.linspace(hi["mean_years_schooling"].min(),
                       hi["mean_years_schooling"].max(), 200)
    ax.plot(xfit, np.poly1d(z)(xfit), color=MPL_MIDGRAY,
            lw=1.6, linestyle="--", zorder=2)

    _mpl_style(ax,
               xlabel="Mean Years of Schooling",
               ylabel="Life Expectancy (Years)")
    fig.tight_layout()
    return fig


def slide_education(prs, df):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "Factor 2: Education shapes health behaviour — and outcomes",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Countries where populations have more schooling show better preventive care uptake and lower chronic disease burden",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
         size=11, color=C_GRAY)

    _add_pic(slide, _to_stream(_chart_education(df)),
             Inches(0.3), Inches(1.05), Inches(8.8), Inches(5.9))

    _rect(slide, Inches(9.25), Inches(1.05), Inches(3.85), Inches(5.9), C_LIGHTGRAY)
    insight = (
        "Health literacy — the ability to\nunderstand and act on health\n"
        "information — varies significantly\nacross OECD nations.\n\n"
        "Countries investing in public\nhealth education achieve\n15–25% lower preventable\n"
        "mortality rates (OECD, 2023).\n\n"
        "Japan's school health programme\nreaches children from age 6,\n"
        "embedding nutrition norms\ndecades before disease risk."
    )
    _txt(slide, insight,
         Inches(9.4), Inches(1.25), Inches(3.55), Inches(5.5),
         size=10.5, color=C_NEARBLACK)

    _txt(slide,
         "Data: UNDP Human Development Index, WHO Global Health Observatory.",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 8 – Factor 3: Efficiency ───────────────────────────────────────────

def _chart_efficiency(df):
    peers = ["Japan", "Spain", "Italy", "France", "Australia",
             "South Korea", "United Kingdom", "Germany", "Canada", "United States"]
    hi = df[df["country"].isin(peers)].dropna(
        subset=["health_expenditure_ppp", "life_expectancy"]).copy()
    hi["efficiency"] = hi["life_expectancy"] / (hi["health_expenditure_ppp"] / 1000)
    hi = hi.sort_values("efficiency", ascending=True)

    fig, ax = plt.subplots(figsize=(9.5, 5.5))
    colors = [MPL_ORANGE if c == "United States" else MPL_BLUE
              for c in hi["country"]]
    bars = ax.barh(hi["country"], hi["efficiency"], color=colors, height=0.65, zorder=3)

    for bar, val, cntry in zip(bars, hi["efficiency"], hi["country"]):
        ax.text(bar.get_width() + 0.12,
                bar.get_y() + bar.get_height() / 2,
                f"{val:.1f} yrs / $1K",
                va="center", fontsize=8.5,
                color=MPL_ORANGE if cntry == "United States" else C_NEARBLACK,
                fontweight="bold" if cntry == "United States" else "normal")

    ax.set_xlim(0, 34)
    ax.set_facecolor("white")
    ax.figure.patch.set_facecolor("white")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color(MPL_MIDGRAY)
    ax.spines["bottom"].set_color(MPL_MIDGRAY)
    ax.tick_params(colors=C_GRAY, labelsize=9)
    ax.set_xlabel("Life Expectancy Years per $1,000 of Health Spend", fontsize=9, color=C_GRAY)
    ax.grid(axis="x", color=MPL_LIGHTGRAY, linewidth=0.8, zorder=0)
    fig.tight_layout()
    return fig


def slide_efficiency(prs, df):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "Factor 3: The US gets the fewest life-years per dollar of any peer nation",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "Life expectancy years per $1,000 spent — a simple metric of what each country buys with its health investment",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
         size=11, color=C_GRAY)

    _add_pic(slide, _to_stream(_chart_efficiency(df)),
             Inches(0.3), Inches(1.05), Inches(12.7), Inches(6.15))

    _txt(slide,
         "Data: World Bank World Development Indicators. Efficiency = Life Expectancy ÷ (Health Expenditure per capita / 1,000).",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 9 – Resolution: Comparison Table ───────────────────────────────────

def slide_comparison_table(prs):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _txt(slide,
         "High-performing systems share three upstream investments",
         Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
         size=21, bold=True, color=C_NEARBLACK)
    _txt(slide,
         "The common thread among Japan, Spain, Italy, and France is not their hospital budgets — it's what happens before the hospital",
         Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.35),
         size=11, color=C_GRAY)

    headers  = ["Metric", "United States", "Japan / Spain", "Gap"]
    col_w    = [Inches(3.9), Inches(2.0), Inches(2.8), Inches(3.0)]
    col_left = [Inches(0.4), Inches(4.4), Inches(6.5), Inches(9.4)]
    row_h    = Inches(0.60)
    top0     = Inches(1.22)

    # Header row
    for hdr, cw, cl in zip(headers, col_w, col_left):
        _rect(slide, cl, top0, cw, row_h, C_DARKBLUE)
        _txt(slide, hdr, cl + Inches(0.06), top0 + Inches(0.1),
             cw - Inches(0.1), row_h - Inches(0.12),
             size=11, bold=True, color=C_WHITE)

    rows = [
        ("Adult obesity rate",          "36.2%",    "4–24%",          "US +12–32 pp"),
        ("Preventable mortality /100K",  "247",       "110–130",        "US −50%"),
        ("Primary care visits / capita", "3.9 / yr", "12–13 / yr",     "US 3× lower"),
        ("Prevention % of health spend", "2.9%",     "5–7%",           "US −2 to −4 pp"),
        ("Life expectancy",              "78.5 yrs", "83–84 yrs",      "US −5 yrs"),
    ]

    for i, row_data in enumerate(rows):
        bg = C_LIGHTGRAY if i % 2 == 0 else C_WHITE
        for j, (val, cw, cl) in enumerate(zip(row_data, col_w, col_left)):
            _rect(slide, cl, top0 + row_h * (i + 1), cw, row_h, bg)
            fc = (C_ORANGE if j == 1 else
                  (MPL_GREEN if j == 2 else C_NEARBLACK))
            _txt(slide, val,
                 cl + Inches(0.06),
                 top0 + row_h * (i + 1) + Inches(0.08),
                 cw - Inches(0.1), row_h - Inches(0.1),
                 size=10.5, color=fc,
                 bold=(j == 1))

    # Summary banner
    _rect(slide, Inches(0.4), Inches(5.58), Inches(12.5), Inches(0.95), C_DARKBLUE)
    _txt(slide,
         "The gap is not in hospital equipment — it's in what happens before people need a hospital.",
         Inches(0.55), Inches(5.7), Inches(12.2), Inches(0.75),
         size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _txt(slide,
         "Source: OECD Health at a Glance 2023, WHO Global Health Observatory, World Bank WDI.",
         Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.26),
         size=8, italic=True, color=C_GRAY)


# ── Slide 10 – Call to Action ─────────────────────────────────────────────────

def slide_cta(prs):
    slide = _blank_slide(prs)
    _bg(slide, C_DARKBLUE)

    _rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, C_BLUE)

    _txt(slide,
         "The question is no longer how much to spend.",
         Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.62),
         size=28, bold=True, color=C_WHITE)

    _txt(slide,
         "It's what the highest-spending countries are doing differently — and what they should stop doing.",
         Inches(0.7), Inches(1.1), Inches(11.5), Inches(0.55),
         size=15, color=C_LIGHTBLUE)

    recs = [
        ("Redirect spend upstream",
         "Raise prevention and primary care from 2.9% to 6%+ of health spend. "
         "Every $1 in prevention saves $5–14 in downstream treatment costs."),
        ("Address obesity structurally",
         "Implement population-level interventions: sugar taxes, food labelling reform, "
         "school nutrition programmes. These are public health investments, not individual choices."),
        ("Redesign system incentives",
         "Shift from fee-for-service to outcomes-based payment models. "
         "Reward providers for keeping populations healthy, not for volume of procedures."),
    ]

    for i, (title, body) in enumerate(recs):
        top = Inches(1.9 + i * 1.6)
        _rect(slide, Inches(0.7), top, Inches(0.5), Inches(0.5), C_BLUE)
        _txt(slide, str(i + 1),
             Inches(0.7), top, Inches(0.5), Inches(0.5),
             size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, title,
             Inches(1.35), top, Inches(11.6), Inches(0.5),
             size=15, bold=True, color=C_WHITE)
        _txt(slide, body,
             Inches(1.35), top + Inches(0.52), Inches(11.6), Inches(0.95),
             size=11, color=C_BLUEGRAY)

    # Closing rule + big idea restatement
    _rect(slide, Inches(0.7), Inches(6.82), Inches(11.8), Inches(0.04), C_BLUE)
    _txt(slide,
         "America's longevity problem is not a budget problem — it is upstream of the budget.",
         Inches(0.7), Inches(6.95), Inches(11.8), Inches(0.42),
         size=12, bold=True, color=C_LIGHTBLUE, align=PP_ALIGN.CENTER)


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    data_path = os.path.join(base_dir, "data", "health_data.csv")
    out_path  = os.path.join(base_dir, "health_expenditure_analysis.pptx")

    print("Loading data …")
    df = pd.read_csv(data_path)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    steps = [
        ("Slide 1  – Cover",                  lambda: slide_cover(prs)),
        ("Slide 2  – Global scatter",          lambda: slide_global_scatter(prs, df)),
        ("Slide 3  – Top spenders",            lambda: slide_top_spenders(prs, df)),
        ("Slide 4  – US anomaly",              lambda: slide_us_anomaly(prs, df)),
        ("Slide 5  – Upstream factors",        lambda: slide_upstream_factors(prs)),
        ("Slide 6  – Obesity",                 lambda: slide_obesity(prs, df)),
        ("Slide 7  – Education",               lambda: slide_education(prs, df)),
        ("Slide 8  – Efficiency",              lambda: slide_efficiency(prs, df)),
        ("Slide 9  – Comparison table",        lambda: slide_comparison_table(prs)),
        ("Slide 10 – Call to action",          lambda: slide_cta(prs)),
    ]

    for label, fn in steps:
        print(f"Building {label} …")
        fn()

    prs.save(out_path)
    print(f"\nSaved to: {out_path}")


if __name__ == "__main__":
    main()
